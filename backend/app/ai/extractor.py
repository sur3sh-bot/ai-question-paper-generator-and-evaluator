"""
Text extraction from uploaded study materials.

Supports: PDF, DOCX, PPTX/PPT, TXT
Libraries: PyMuPDF (fitz), python-docx, python-pptx
"""

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_pdf_text(file_path: str) -> str:
    """
    Extract text from a PDF file using PyMuPDF.

    Iterates over all pages, skips blank ones, and joins the content.

    Args:
        file_path: Absolute or relative path to the PDF file.

    Returns:
        Concatenated text from all non-empty pages.

    Raises:
        RuntimeError: If the file cannot be opened or read.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("PyMuPDF is not installed. Run: pip install pymupdf")

    try:
        doc = fitz.open(file_path)
        pages_text = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text and text.strip():
                pages_text.append(text)
            else:
                logger.debug("PDF page %d is empty, skipping.", page_num)
        doc.close()
        return "\n\n".join(pages_text)
    except Exception as exc:
        raise RuntimeError(f"Failed to extract text from PDF '{file_path}': {exc}") from exc


def extract_docx_text(file_path: str) -> str:
    """
    Extract text from a DOCX file using python-docx.

    Reads all paragraphs and table cells, preserving paragraph breaks.

    Args:
        file_path: Absolute or relative path to the DOCX file.

    Returns:
        Plain text with paragraphs separated by newlines.

    Raises:
        RuntimeError: If the file cannot be opened or read.
    """
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    try:
        doc = Document(file_path)
        parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append(" | ".join(row_texts))

        return "\n\n".join(parts)
    except Exception as exc:
        raise RuntimeError(f"Failed to extract text from DOCX '{file_path}': {exc}") from exc


def extract_ppt_text(file_path: str) -> str:
    """
    Extract text from a PPTX file using python-pptx.

    Reads all slides, iterates over text frames and table cells.
    Empty slides are skipped.

    Args:
        file_path: Absolute or relative path to the PPTX file.

    Returns:
        Text from all slides, with slide separators.

    Raises:
        RuntimeError: If the file cannot be opened or read.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    try:
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts = []

            for shape in slide.shapes:
                # Text frames (titles, content boxes)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)

                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_parts.append(" | ".join(row_texts))

            if slide_parts:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
            else:
                logger.debug("Slide %d is empty, skipping.", slide_num)

        return "\n\n".join(slides_text)
    except Exception as exc:
        raise RuntimeError(f"Failed to extract text from PPTX '{file_path}': {exc}") from exc


def extract_txt_text(file_path: str) -> str:
    """
    Read a plain text file with automatic encoding detection.

    Tries UTF-8 first, falls back to latin-1 to handle legacy files.

    Args:
        file_path: Absolute or relative path to the TXT file.

    Returns:
        Full file content as a string.

    Raises:
        RuntimeError: If the file cannot be read.
    """
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as fh:
                return fh.read()
        except UnicodeDecodeError:
            continue
        except Exception as exc:
            raise RuntimeError(f"Failed to read TXT file '{file_path}': {exc}") from exc

    raise RuntimeError(f"Could not decode '{file_path}' with any supported encoding.")


def extract_text(file_path: str, extension: str) -> str:
    """
    Dispatcher: route to the correct extractor based on file extension.

    Args:
        file_path: Path to the uploaded file.
        extension: Lowercase extension without dot (e.g., 'pdf', 'docx').

    Returns:
        Extracted text string.

    Raises:
        ValueError: If the extension is unsupported.
        RuntimeError: If extraction fails.
    """
    dispatch = {
        "pdf":  extract_pdf_text,
        "docx": extract_docx_text,
        "pptx": extract_ppt_text,
        "ppt":  extract_ppt_text,
        "txt":  extract_txt_text,
    }

    handler = dispatch.get(extension.lower())
    if not handler:
        raise ValueError(
            f"Unsupported file type '.{extension}'. "
            f"Supported types: {', '.join(dispatch.keys())}"
        )

    logger.info("Extracting text from '%s' (type=%s)", file_path, extension)
    text = handler(file_path)

    if not text or not text.strip():
        raise ValueError("The uploaded file appears to be empty or contains no extractable text.")

    logger.info("Extracted %d characters from '%s'.", len(text), file_path)
    return text